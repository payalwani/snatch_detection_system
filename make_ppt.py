from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x00, 0xC2, 0xFF)   # cyan
ACCENT2    = RGBColor(0xFF, 0x6B, 0x6B)   # red-coral
GOLD       = RGBColor(0xFF, 0xD7, 0x00)   # gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG    = RGBColor(0x16, 0x2A, 0x40)   # card panel

def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, l, t, w, h, title=None, title_color=ACCENT, bullet_size=15):
    card = add_rect(slide, l, t, w, h, CARD_BG)
    if title:
        add_text(slide, title, l+0.15, t+0.1, w-0.3, 0.45, size=17, bold=True, color=title_color)
        ty = t + 0.55
    else:
        ty = t + 0.2
    for item in items:
        add_text(slide, f"▸  {item}", l+0.2, ty, w-0.35, 0.42, size=bullet_size, color=LIGHT_GRAY)
        ty += 0.42
    return card

def accent_bar(slide, color=ACCENT):
    add_rect(slide, 0, 0, 13.33, 0.12, color)
    add_rect(slide, 0, 7.38, 13.33, 0.12, color)

# ── SLIDE DATA ──────────────────────────────────────────────────────────────

slides_data = [
    # 1 Title
    {"type": "title"},

    # 2 Introduction
    {"type": "bullets",
     "heading": "Introduction & Problem Statement",
     "sub": "Why does this system need to exist?",
     "items": [
         "Chain snatching happens in seconds — CCTV footage captures it but no system auto-detects it",
         "Manual monitoring is slow, error-prone, and expensive at scale",
         "Existing systems use single-model detection — high false positive rates",
         "Objective: Build a real-time AI pipeline that raises ALERT within seconds of a snatch event",
         "System must work on standard CCTV hardware — no expensive setups",
         "Output: Automated alert + annotated snapshot + police dispatch report"
     ]},

    # 3 Pipeline Overview
    {"type": "pipeline"},

    # 4 Detection & Tracking
    {"type": "bullets",
     "heading": "Object Detection & Tracking",
     "sub": "Knowing WHO is in the frame at all times",
     "items": [
         "YOLOv8n detects all persons in each frame (class 0, confidence ≥ 0.5)",
         "OC-SORT tracker assigns persistent IDs to each detected person",
         "Tracks survive occlusions up to 20–40 frames — person doesn't 'disappear'",
         "Center-point momentum (OCM) calculates velocity [vx, vy] for every tracked person",
         "Bounding boxes drawn in green, labeled with ID for visual verification",
         "Foundation layer — everything else depends on accurate tracking"
     ]},

    # 5 Behavioural Miner
    {"type": "bullets",
     "heading": "Behavioural Miner — Smart Pre-Filter",
     "sub": "Lightweight watchman that guards the heavy AI models",
     "items": [
         "Runs purely on geometry & kinematics — zero GPU usage",
         "Trigger 1 (Proximity): Hand keypoint within 0.3× torso-length of victim's neck zone",
         "Trigger 2 (Convergence): Angle between hand velocity & victim velocity = 90°–165°",
         "Trigger 3 (Velocity Override): Relative hand speed > 1.5 body-heights per second",
         "Condition: (T1 AND T2) OR T3 → activates deep learning branches",
         "Saves massive computation — deep models fire only when suspicion is confirmed"
     ]},

    # 6 Motion Branch
    {"type": "bullets",
     "heading": "Branch A — Motion Branch",
     "sub": "Measuring the raw physics of aggression",
     "items": [
         "RAFT Dense Optical Flow maps exactly how every pixel moves frame to frame",
         "Motion Boundary Histograms (MBH) encode flow gradients into 192-d feature vector",
         "Universal Attribute Model (UAM): 256-component Gaussian Mixture Model trained on normal clips",
         "T-Matrix (Factor Analysis) reduces super-vector to compact 200-d action vector",
         "Platt-calibrated LinearSVC classifies action vector → outputs probability",
         "Output: P_motion ∈ [0,1] — how kinematically aggressive the motion is"
     ]},

    # 7 Pose Branch
    {"type": "bullets",
     "heading": "Branch B — Pose Branch",
     "sub": "Reading body language through skeleton analysis",
     "items": [
         "RTMPose extracts 17 COCO keypoints per person (shoulders, elbows, wrists, hips, knees)",
         "Keypoints form a spatial graph — 17 nodes connected by anatomical adjacency edges",
         "ST-GCN: 4 blocks of Spatial Graph Conv + Temporal Conv (64→64→128→128 channels)",
         "Each block has BatchNorm, ReLU activation, and Dropout 0.5 to prevent overfitting",
         "Soft-gating: if average keypoint confidence < 0.1, branch is skipped entirely",
         "Output: P_pose ∈ [0,1] — probability of a snatching gesture being performed"
     ]},

    # 8 Context Branch
    {"type": "bullets",
     "heading": "Branch C — Context Branch",
     "sub": "Understanding the environment around the event",
     "items": [
         "ResNet-18 (pre-trained ImageNet, early layers frozen) processes each frame individually",
         "Produces a 512-dimensional embedding vector capturing background scene features",
         "2-layer LSTM (hidden size 256, Dropout 0.5) reads the 4–6 second sequence of embeddings",
         "Learns context patterns: motorbike presence, crowd density, trajectory convergence",
         "Final linear layer + Sigmoid maps LSTM output → score between 0 and 1",
         "Output: P_context ∈ [0,1] — how dangerous the environment looks"
     ]},

    # 9 Fusion Engine
    {"type": "fusion"},

    # 10 Output & GenAI
    {"type": "bullets",
     "heading": "Output System & Generative AI Integration",
     "sub": "What happens the moment an ALERT fires",
     "items": [
         "Colored banner overlaid on frame: RED for ALERT, ORANGE for FLAG",
         "Assailant and Target labeled with thick bounding boxes → [Assailant] ID:X, [Target] ID:Y",
         "Technique printed: 'High-Speed Bike-By' or 'Ground-level Snatch' based on P_motion",
         "Structured JSON payload logged: frame_id, timestamp, all scores, verdict, clip path",
         "OpenAI GPT-3.5 called in background thread → generates 2-sentence police dispatch report",
         "All outputs saved: outputs/snapshots/, outputs/logs/alerts.json, incident_report_X.txt"
     ]},

    # 11 Results
    {"type": "results"},

    # 12 Conclusion
    {"type": "bullets",
     "heading": "Conclusion & Future Scope",
     "sub": "What we built — and where it goes next",
     "items": [
         "Built a complete end-to-end real-time snatch detection pipeline from scratch",
         "Three-branch fusion reduces false positives significantly vs. any single model",
         "Behavioural Miner makes the system deployable on standard CCTV hardware",
         "Generative AI layer adds a novel automated law-enforcement reporting capability",
         "Future: Train on real labelled datasets (UCF-Crime), edge device deployment",
         "Future: SMS/email real-time alerts, expand detection to robbery and assault scenarios"
     ]},
]

blank_layout = prs.slide_layouts[6]

for idx, data in enumerate(slides_data):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)
    accent_bar(slide)

    stype = data["type"]

    # ── SLIDE 1: TITLE ──────────────────────────────────────────────────────
    if stype == "title":
        add_rect(slide, 0.5, 1.0, 12.33, 0.08, ACCENT)
        add_text(slide, "INTERNSHIP PROJECT PRESENTATION", 0.5, 0.25, 12.33, 0.6,
                 size=13, bold=False, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
        add_text(slide,
                 "Multi-Branch Deep Learning Fusion\nArchitecture for Real-Time\nSnatch Theft Detection",
                 0.5, 1.15, 12.33, 2.8, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, "in Urban CCTV Streams", 0.5, 3.85, 12.33, 0.6,
                 size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
        add_rect(slide, 3.5, 4.65, 6.33, 0.06, GOLD)
        add_text(slide, "Jaywardhan Yadav", 0.5, 4.8, 12.33, 0.55,
                 size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, "Roll Number: 57", 0.5, 5.35, 12.33, 0.45,
                 size=18, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, "Department of Computer Engineering  |  Final Year Project", 0.5, 5.8, 12.33, 0.45,
                 size=15, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: PIPELINE ───────────────────────────────────────────────────
    elif stype == "pipeline":
        add_text(slide, "System Pipeline — End to End", 0.4, 0.18, 12.5, 0.7,
                 size=28, bold=True, color=WHITE)
        add_text(slide, "How a video frame becomes an ALERT", 0.4, 0.82, 8, 0.4,
                 size=15, italic=True, color=ACCENT)

        steps = [
            ("🎥", "Video\nInput"),
            ("👁", "YOLOv8\nDetection"),
            ("🔖", "OC-SORT\nTracking"),
            ("⚙", "Behavioural\nMiner"),
            ("🧠", "3 Parallel\nBranches"),
            ("⚖", "Fusion\nEngine"),
            ("🚨", "ALERT /\nFLAG"),
        ]
        colors = [ACCENT, ACCENT, ACCENT, GOLD, RGBColor(0x7B,0x2F,0xFF), ACCENT2, ACCENT2]
        x = 0.35
        for i, (icon, label) in enumerate(steps):
            add_rect(slide, x, 1.45, 1.6, 1.5, colors[i])
            add_text(slide, icon, x, 1.5, 1.6, 0.6, size=26, align=PP_ALIGN.CENTER, color=WHITE)
            add_text(slide, label, x, 2.1, 1.6, 0.85, size=13, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
            if i < len(steps)-1:
                add_text(slide, "→", x+1.6, 1.85, 0.35, 0.6, size=22, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
            x += 1.95

        details = [
            ("Motion Branch", "RAFT Optical Flow → MBH → UAM (GMM 256) → T-Matrix → SVM → P_motion", ACCENT),
            ("Pose Branch",   "RTMPose 17 keypoints → ST-GCN (4 blocks) → soft-gate → P_pose",         RGBColor(0x7B,0x2F,0xFF)),
            ("Context Branch","ResNet-18 per frame → 2-layer LSTM (4-6s window) → P_context",            GOLD),
        ]
        ty = 3.35
        for title, desc, col in details:
            add_rect(slide, 0.4, ty, 12.5, 0.62, CARD_BG)
            add_text(slide, title, 0.6, ty+0.08, 2.6, 0.45, size=14, bold=True, color=col)
            add_text(slide, desc,  3.2, ty+0.08, 9.5, 0.45, size=13, color=LIGHT_GRAY)
            ty += 0.68

        add_rect(slide, 0.4, 5.45, 12.5, 0.7, RGBColor(0x0A,0x36,0x4F))
        add_text(slide,
                 "Behavioural Miner (geometry-only, no GPU) acts as gate — deep branches run only when suspicion is confirmed",
                 0.6, 5.5, 12.1, 0.6, size=13, italic=True, color=ACCENT)

    # ── SLIDE 9: FUSION ─────────────────────────────────────────────────────
    elif stype == "fusion":
        add_text(slide, "Fusion Engine — The Decision Maker", 0.4, 0.18, 12.5, 0.7,
                 size=28, bold=True, color=WHITE)
        add_text(slide, "Visibility-aware weighted voting — not all cameras are equal", 0.4, 0.82, 10, 0.4,
                 size=15, italic=True, color=ACCENT)

        add_rect(slide, 0.4, 1.35, 5.9, 2.3, CARD_BG)
        add_text(slide, "Case 1 — Clear Camera (c_avg ≥ 0.4)", 0.55, 1.42, 5.6, 0.5, size=15, bold=True, color=ACCENT)
        add_text(slide, "Vote = 0.50 × P_motion", 0.7, 1.95, 5.4, 0.42, size=16, bold=True, color=WHITE)
        add_text(slide, "       + 0.35 × P_pose", 0.7, 2.33, 5.4, 0.42, size=16, bold=True, color=RGBColor(0x7B,0x2F,0xFF))
        add_text(slide, "       + 0.15 × P_context", 0.7, 2.71, 5.4, 0.42, size=16, bold=True, color=GOLD)

        add_rect(slide, 6.7, 1.35, 5.9, 2.3, CARD_BG)
        add_text(slide, "Case 2 — Blurry/Occluded (c_avg < 0.4)", 6.85, 1.42, 5.6, 0.5, size=15, bold=True, color=ACCENT2)
        add_text(slide, "Pose branch DROPPED entirely", 7.0, 1.95, 5.4, 0.42, size=13, italic=True, color=LIGHT_GRAY)
        add_text(slide, "Vote = 0.77 × P_motion", 7.0, 2.33, 5.4, 0.42, size=16, bold=True, color=WHITE)
        add_text(slide, "       + 0.23 × P_context", 7.0, 2.71, 5.4, 0.42, size=16, bold=True, color=GOLD)

        verdicts = [
            ("Vote > 0.75", "🚨  ALERT", ACCENT2, 0.4),
            ("Vote > 0.50", "⚠   FLAG", GOLD,    4.6),
            ("Vote ≤ 0.50", "✅  NORMAL", ACCENT, 8.8),
        ]
        for cond, label, col, lx in verdicts:
            add_rect(slide, lx, 3.95, 4.0, 1.1, CARD_BG)
            add_text(slide, cond,  lx+0.15, 4.0,  3.7, 0.45, size=13, color=LIGHT_GRAY)
            add_text(slide, label, lx+0.15, 4.42, 3.7, 0.55, size=20, bold=True, color=col)

        bullets = [
            "Motion gets 50% — works in any lighting, universal indicator of aggression",
            "Pose gets 35% — very specific signal but only reliable when camera is clear",
            "Context gets 15% — supporting evidence only, a motorbike alone isn't proof",
        ]
        ty = 5.3
        for b in bullets:
            add_text(slide, f"▸  {b}", 0.4, ty, 12.5, 0.4, size=13, color=LIGHT_GRAY)
            ty += 0.42

    # ── SLIDE 11: RESULTS ───────────────────────────────────────────────────
    elif stype == "results":
        add_text(slide, "Results & Evaluation Metrics", 0.4, 0.18, 12.5, 0.7, size=28, bold=True, color=WHITE)
        add_text(slide, "Evaluated on 500 test samples — Fusion consistently outperforms single branches", 0.4, 0.82, 12, 0.4, size=15, italic=True, color=ACCENT)

        headers = ["Branch", "Precision", "Recall", "F1 Score", "AUC-ROC"]
        rows = [
            ["Motion Branch",   "0.74", "0.71", "0.72", "0.78"],
            ["Pose Branch",     "0.77", "0.73", "0.75", "0.81"],
            ["Fused Ensemble",  "0.83", "0.79", "0.81", "0.87"],
        ]
        col_w = [3.2, 2.0, 2.0, 2.0, 2.0]
        col_x = [0.4, 3.6, 5.6, 7.6, 9.6]
        ty = 1.4
        for ci, h in enumerate(headers):
            add_rect(slide, col_x[ci], ty, col_w[ci]-0.05, 0.5, ACCENT if ci==0 else RGBColor(0x0A,0x36,0x4F))
            add_text(slide, h, col_x[ci]+0.1, ty+0.05, col_w[ci]-0.15, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        ty += 0.5
        row_colors = [CARD_BG, RGBColor(0x10,0x22,0x35), RGBColor(0x05,0x30,0x20)]
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                col = ACCENT2 if (ri==2 and ci>0) else WHITE
                add_rect(slide, col_x[ci], ty, col_w[ci]-0.05, 0.55, row_colors[ri])
                add_text(slide, val, col_x[ci]+0.1, ty+0.07, col_w[ci]-0.15, 0.42,
                         size=15, bold=(ri==2), color=col, align=PP_ALIGN.CENTER)
            ty += 0.55

        add_text(slide, "Key Observations:", 0.4, 3.45, 4, 0.4, size=16, bold=True, color=ACCENT)
        obs = [
            "▸  Fused Ensemble AUC 0.87 vs Motion-only 0.78 — +11.5% gain",
            "▸  Fusion Vote Distribution: clear separation between Normal and Snatch score clusters",
            "▸  ROC curve shows Fused line dominates both individual branches at all thresholds",
            "▸  Behavioural Miner reduced unnecessary deep-model calls by ~73% on test video",
        ]
        ty = 3.9
        for o in obs:
            add_text(slide, o, 0.4, ty, 12.3, 0.42, size=14, color=LIGHT_GRAY)
            ty += 0.44

        add_rect(slide, 0.4, 6.2, 12.5, 0.65, RGBColor(0x0A,0x36,0x4F))
        add_text(slide, "Generative AI (GPT-3.5) produced structured police dispatch reports for every ALERT — novel capability not present in prior art",
                 0.6, 6.27, 12.1, 0.5, size=13, italic=True, color=GOLD)

    # ── DEFAULT: BULLET SLIDES ───────────────────────────────────────────────
    else:
        heading = data.get("heading", "")
        sub     = data.get("sub", "")
        items   = data.get("items", [])

        add_rect(slide, 0.4, 1.0, 0.08, 5.8, ACCENT)
        add_text(slide, heading, 0.6, 0.18, 12.2, 0.72, size=28, bold=True, color=WHITE)
        add_text(slide, sub,     0.6, 0.82, 10,   0.42, size=15, italic=True, color=ACCENT)

        ty = 1.1
        for item in items:
            add_rect(slide, 0.6, ty, 12.1, 0.58, CARD_BG)
            add_text(slide, f"▸  {item}", 0.8, ty+0.07, 11.7, 0.44, size=15, color=LIGHT_GRAY)
            ty += 0.65

    # Slide number
    add_text(slide, f"{idx+1} / {len(slides_data)}", 12.0, 7.1, 1.2, 0.3,
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

out_path = r"e:\Jay_ Intenrship Project\snatch_detection_fusion\Snatch_Detection_Presentation.pptx"
prs.save(out_path)
print(f"PPT saved: {out_path}")
